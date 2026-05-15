import os
import time

class DocumentFormats:
    @staticmethod
    def convertir(origen_path: str, destino_path: str, from_ext: str, to_ext: str) -> str:
        """
        Enruta la conversión a la función correspondiente según el formato origen y destino.
        """
        origen_path = os.path.abspath(origen_path)
        destino_path = os.path.abspath(destino_path)
        
        try:
            if from_ext == "jpg" and to_ext == "pdf":
                return DocumentFormats.jpg_to_pdf(origen_path, destino_path)
            elif from_ext == "pdf" and to_ext == "jpg":
                return DocumentFormats.pdf_to_jpg(origen_path, destino_path)
            elif from_ext == "pdf" and to_ext == "docx":
                return DocumentFormats.pdf_to_word(origen_path, destino_path)
            elif from_ext == "docx" and to_ext == "pdf":
                return DocumentFormats.word_to_pdf(origen_path, destino_path)
            elif from_ext == "xlsx" and to_ext == "pdf":
                return DocumentFormats.excel_to_pdf(origen_path, destino_path)
            elif from_ext == "pptx" and to_ext == "pdf":
                return DocumentFormats.powerpoint_to_pdf(origen_path, destino_path)
            elif from_ext == "pdf" and to_ext == "xlsx":
                return DocumentFormats.pdf_to_excel(origen_path, destino_path)
            elif from_ext == "pdf" and to_ext == "pptx":
                return DocumentFormats.pdf_to_powerpoint(origen_path, destino_path)
            elif from_ext == "html" and to_ext == "pdf":
                return DocumentFormats.html_to_pdf(origen_path, destino_path)
            else:
                return f"❌ Conversión de {from_ext.upper()} a {to_ext.upper()} no está soportada aún."
        except Exception as e:
            return f"❌ Error al convertir: {e}"

    @staticmethod
    def jpg_to_pdf(origen, destino):
        import img2pdf
        from PIL import Image
        try:
            # Validate if it's an image
            with Image.open(origen) as img:
                img.verify()
                
            with open(destino, "wb") as f:
                f.write(img2pdf.convert(origen))
            return "✅ Archivo JPG convertido a PDF exitosamente."
        except Exception as e:
            return f"❌ Error al convertir JPG a PDF: {e}"

    @staticmethod
    def pdf_to_jpg(origen, destino):
        import fitz
        import os
        
        doc = fitz.open(origen)
        num_pages = len(doc)
        
        if num_pages == 0:
            return "❌ El PDF está vacío."
            
        output_dir = os.path.dirname(destino)
        base_name = os.path.splitext(os.path.basename(destino))[0]
        
        # Si tiene más de 10 hojas, creamos una carpeta dedicada
        if num_pages > 5:
            target_dir = os.path.join(output_dir, base_name)
            if not os.path.exists(target_dir):
                os.makedirs(target_dir)
        else:
            target_dir = output_dir
            
        # Extraer todas las páginas
        for i in range(num_pages):
            page = doc.load_page(i)
            pix = page.get_pixmap(dpi=300)
            
            # Si es solo 1 hoja, mantiene el nombre original, sino enumera
            if num_pages == 1:
                file_path = os.path.join(target_dir, f"{base_name}.jpg")
            else:
                file_path = os.path.join(target_dir, f"{base_name}_pagina_{i+1}.jpg")
                
            pix.save(file_path)
            
        doc.close()
        
        if num_pages > 5:
            return f"✅ {num_pages} páginas extraídas a JPG y guardadas en la carpeta '{base_name}'."
        elif num_pages > 1:
            return f"✅ {num_pages} páginas extraídas a JPG exitosamente."
        else:
            return "✅ Página convertida a JPG exitosamente."

    @staticmethod
    def pdf_to_word(origen, destino):
        from pdf2docx import Converter
        cv = None
        try:
            cv = Converter(origen)
            cv.convert(destino, start=0, end=None)
            
            # Post-procesamiento para arreglar el bug de salto de página
            try:
                import docx
                from docx.shared import Cm
                doc = docx.Document(destino)
                for section in doc.sections:
                    # Reducir los márgenes a 0.5 cm para dar espacio extra y evitar saltos de página
                    section.top_margin = Cm(0.5)
                    section.bottom_margin = Cm(0.5)
                    section.left_margin = Cm(0.5)
                    section.right_margin = Cm(0.5)
                doc.save(destino)
            except Exception as e_margin:
                print(f"No se pudieron ajustar los márgenes: {e_margin}")

            return "✅ Archivo PDF convertido a Word exitosamente (Márgenes optimizados)."
        except Exception as e:
            return f"❌ Error al convertir PDF a Word: {e}"
        finally:
            if cv is not None:
                cv.close()

    @staticmethod
    def word_to_pdf(origen, destino):
        try:
            from docx2pdf import convert
            convert(origen, destino)
            return "✅ Archivo Word convertido a PDF exitosamente."
        except Exception as e:
            # Fallback a comtypes
            try:
                import comtypes.client
                word = comtypes.client.CreateObject("Word.Application")
                try:
                    word.Visible = False
                    doc = word.Documents.Open(origen)
                    doc.SaveAs(destino, FileFormat=17) # 17 = wdFormatPDF
                    doc.Close()
                    return "✅ Archivo Word convertido a PDF exitosamente."
                finally:
                    word.Quit()
            except Exception as e2:
                return f"❌ Error al convertir Word a PDF: {e} | Fallback COM: {e2}"

    @staticmethod
    def excel_to_pdf(origen, destino):
        import comtypes.client
        excel = None
        wb = None
        try:
            excel = comtypes.client.CreateObject("Excel.Application")
            excel.Visible = False
            wb = excel.Workbooks.Open(origen)
            # 0 = xlTypePDF
            wb.ExportAsFixedFormat(0, destino)
            return "✅ Archivo Excel convertido a PDF exitosamente."
        except Exception as e:
            return f"❌ Error usando Excel (comtypes): {e}. Asegúrate de tener MS Excel instalado."
        finally:
            if wb is not None:
                wb.Close(False)
            if excel is not None:
                excel.Quit()

    @staticmethod
    def powerpoint_to_pdf(origen, destino):
        import comtypes.client
        powerpoint = None
        presentation = None
        try:
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            # 32 = ppFixedFormatTypePDF
            presentation = powerpoint.Presentations.Open(origen, WithWindow=False)
            presentation.ExportAsFixedFormat(destino, 32)
            return "✅ Archivo PowerPoint convertido a PDF exitosamente."
        except Exception as e:
            return f"❌ Error usando PowerPoint (comtypes): {e}. Asegúrate de tener MS PowerPoint instalado."
        finally:
            if presentation is not None:
                presentation.Close()
            if powerpoint is not None:
                powerpoint.Quit()

    @staticmethod
    def pdf_to_excel(origen, destino):
        import pdfplumber
        import pandas as pd
        try:
            with pdfplumber.open(origen) as pdf:
                all_tables = []
                for page in pdf.pages:
                    tables = page.extract_tables()
                    for table in tables:
                        if not table or not table[0]:
                            continue
                        
                        # Limpiar las cabeceras (primera fila) para evitar Nones y duplicados
                        headers = table[0]
                        clean_headers = []
                        for i, h in enumerate(headers):
                            val = str(h).strip() if h else f"Columna_{i+1}"
                            # Asegurar que no hay duplicados en el nombre
                            if val in clean_headers:
                                val = f"{val}_{i+1}"
                            clean_headers.append(val)
                            
                        # Limpiar el cuerpo asegurando que coincida el número de columnas
                        data = []
                        for row in table[1:]:
                            clean_row = [str(cell).strip() if cell else "" for cell in row]
                            # Ajustar a la longitud de headers
                            if len(clean_row) < len(clean_headers):
                                clean_row.extend([""] * (len(clean_headers) - len(clean_row)))
                            elif len(clean_row) > len(clean_headers):
                                clean_row = clean_row[:len(clean_headers)]
                            data.append(clean_row)

                        df = pd.DataFrame(data, columns=clean_headers)
                        all_tables.append(df)
                
                if all_tables:
                    with pd.ExcelWriter(destino) as writer:
                        for i, df in enumerate(all_tables):
                            df.to_excel(writer, sheet_name=f"Tabla_{i+1}", index=False)
                    return "✅ Tablas del PDF extraídas a Excel exitosamente."
                else:
                    pd.DataFrame(["No se detectaron tablas claras en el PDF"]).to_excel(destino, index=False)
                    return "✅ Archivo Excel creado, pero sin tablas definidas."
        except Exception as e:
            return f"❌ Error extrayendo Excel del PDF: {e}"

    @staticmethod
    def pdf_to_powerpoint(origen, destino):
        import fitz
        from pptx import Presentation
        from pptx.util import Inches
        import tempfile
        import shutil
        import os
        import time
        
        temp_dir = tempfile.mkdtemp(prefix="mediahub_pdf2ppt_")
        try:
            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[6] # Blank
            
            doc = fitz.open(origen)
            for i in range(len(doc)):
                page = doc.load_page(i)
                pix = page.get_pixmap(dpi=150)
                img_path = os.path.join(temp_dir, f"page_{i}.png")
                pix.save(img_path)
                
                slide = prs.slides.add_slide(blank_slide_layout)
                slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            
            prs.save(destino)
            doc.close()
            return "✅ Páginas del PDF convertidas a diapositivas de PowerPoint exitosamente."
        except Exception as e:
            return f"❌ Error al convertir PDF a PowerPoint: {e}"
        finally:
            if os.path.exists(temp_dir):
                shutil.rmtree(temp_dir, ignore_errors=True)

    @staticmethod
    def html_to_pdf(origen, destino):
        from PyQt6.QtGui import QTextDocument
        from PyQt6.QtPrintSupport import QPrinter
        try:
            try:
                with open(origen, "r", encoding="utf-8") as f:
                    html_content = f.read()
            except UnicodeDecodeError:
                with open(origen, "r", encoding="latin-1") as f:
                    html_content = f.read()
                
            doc = QTextDocument()
            doc.setHtml(html_content)
            
            printer = QPrinter(QPrinter.PrinterMode.HighResolution)
            printer.setOutputFormat(QPrinter.OutputFormat.PdfFormat)
            printer.setOutputFileName(destino)
            
            doc.print(printer)
            if os.path.exists(destino):
                return "✅ Archivo HTML convertido a PDF exitosamente."
            else:
                return "❌ No se pudo crear el archivo PDF. Asegúrate de tener una interfaz gráfica disponible."
        except Exception as e:
            return f"❌ Error al convertir HTML: {e}"

